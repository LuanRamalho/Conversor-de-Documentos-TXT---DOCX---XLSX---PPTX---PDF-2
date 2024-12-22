import tkinter as tk
from tkinter import filedialog, messagebox
from PyPDF2 import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation

def choose_file():
    global file_path
    file_path = filedialog.askopenfilename(filetypes=[("PDF Files", "*.pdf")])
    if file_path:
        lbl_file_chosen.config(text=f"Arquivo: {file_path}")

def convert_file():
    if not file_path:
        messagebox.showerror("Erro", "Por favor, escolha um arquivo PDF!")
        return

    output_format = output_option.get()
    if not output_format:
        messagebox.showerror("Erro", "Por favor, escolha um formato de saída!")
        return

    try:
        reader = PdfReader(file_path)
        text_content = ""
        for page in reader.pages:
            text_content += page.extract_text()

        if output_format == "TXT":
            output_file = file_path.replace(".pdf", ".txt")
            with open(output_file, "w", encoding="utf-8") as txt_file:
                txt_file.write(text_content)
        elif output_format == "DOCX":
            output_file = file_path.replace(".pdf", ".docx")
            doc = Document()
            doc.add_paragraph(text_content)
            doc.save(output_file)
        elif output_format == "XLSX":
            output_file = file_path.replace(".pdf", ".xlsx")
            workbook = openpyxl.Workbook()
            sheet = workbook.active
            for i, line in enumerate(text_content.splitlines(), start=1):
                sheet.cell(row=i, column=1, value=line)
            workbook.save(output_file)
        elif output_format == "PPTX":
            output_file = file_path.replace(".pdf", ".pptx")
            presentation = Presentation()
            for line in text_content.splitlines():
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                textbox = slide.shapes.add_textbox(left=0, top=0, width=Inches(10), height=Inches(2))
                frame = textbox.text_frame
                frame.text = line
            presentation.save(output_file)

        messagebox.showinfo("Sucesso", f"Arquivo convertido e salvo como: {output_file}")
    except Exception as e:
        messagebox.showerror("Erro", f"Erro ao converter o arquivo: {e}")

# Interface gráfica
root = tk.Tk()
root.title("Conversor de Arquivos PDF")
root.geometry("400x300")
root.configure(bg="#f0f8ff")

file_path = None

# Caixa de opção de entrada
lbl_input = tk.Label(root, text="Formato de Entrada:", bg="#f0f8ff", font=("Arial", 12))
lbl_input.pack(pady=10)
input_option = tk.StringVar(value="PDF")
input_dropdown = tk.OptionMenu(root, input_option, "PDF")
input_dropdown.configure(state="disabled", bg="#add8e6", font=("Arial", 10))
input_dropdown.pack()

# Caixa de opções de saída
lbl_output = tk.Label(root, text="Formato de Saída:", bg="#f0f8ff", font=("Arial", 12))
lbl_output.pack(pady=10)
output_option = tk.StringVar()
output_dropdown = tk.OptionMenu(root, output_option, "TXT", "DOCX", "XLSX", "PPTX")
output_dropdown.configure(bg="#add8e6", font=("Arial", 10))
output_dropdown.pack()

# Botão para escolher arquivo
btn_choose = tk.Button(root, text="Escolher Arquivo", command=choose_file, bg="#87cefa", font=("Arial", 12))
btn_choose.pack(pady=10)
lbl_file_chosen = tk.Label(root, text="Nenhum arquivo escolhido", bg="#f0f8ff", font=("Arial", 10))
lbl_file_chosen.pack()

# Botão para realizar conversão
btn_convert = tk.Button(root, text="Converter Arquivo", command=convert_file, bg="#4682b4", fg="white", font=("Arial", 12))
btn_convert.pack(pady=20)

root.mainloop()
